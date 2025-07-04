import pandas as pd
import streamlit as st
import re
import os
import gc
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import numpy as np
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go

# Matplotlib imports for PPT generation
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import seaborn as sns

# Set Streamlit page configuration
st.set_page_config(
    layout='wide',
    page_title="📊 Excel Dashboard - Data Table & Visualizations",
    page_icon="📊",
    initial_sidebar_state="expanded"
)

# Add cloud-optimized CSS
def add_cloud_css():
    st.markdown("""
    <style>
        .main .block-container {
            padding-top: 1rem;
            padding-bottom: 1rem;
            max-width: 95%;
        }
        .js-plotly-plot, .plotly {
            width: 100% !important;
            height: auto !important;
        }
        .element-container {
            margin: 0.5rem 0;
        }
        .dataframe {
            font-size: 12px;
        }
        [data-testid="metric-container"] {
            background-color: #f0f2f6;
            border: 1px solid #e0e0e0;
            padding: 0.5rem;
            border-radius: 0.5rem;
            margin: 0.25rem;
        }
        .stTabs [data-baseweb="tab-list"] {
            gap: 4px;
        }
        .stTabs [data-baseweb="tab"] {
            height: 40px;
            padding-left: 15px;
            padding-right: 15px;
            background-color: #f0f2f6;
            border-radius: 8px 8px 0 0;
            font-size: 14px;
        }
        .streamlit-expanderHeader {
            background-color: #f0f2f6;
            border-radius: 5px;
            font-size: 14px;
        }
        @media (max-width: 768px) {
            .main .block-container {
                padding: 0.5rem;
            }
            .stTabs [data-baseweb="tab"] {
                height: 35px;
                padding: 0 10px;
                font-size: 12px;
            }
        }
    </style>
    """, unsafe_allow_html=True)

add_cloud_css()

st.title("📊 Excel Dashboard - Data Table & Visualizations")

# Define exclusion terms for branches
BRANCH_EXCLUDE_TERMS = ['CHN Total', 'ERD SALES', 'North Total', 'WEST SALES', 'GROUP COMPANIES']

# Memory management
def optimize_memory():
    gc.collect()

def is_streamlit_cloud():
    """Check if running on Streamlit Cloud"""
    import os
    return os.getenv('STREAMLIT_SHARING_MODE') is not None or 'streamlit.io' in os.getenv('HOSTNAME', '')

# Utility functions
def safe_convert_value(x):
    try:
        if x is None or (hasattr(x, 'isna') and pd.isna(x)) or pd.isna(x):
            return None
        str_val = str(x)
        if str_val.lower() in ['nan', 'none', 'nat', '', 'null']:
            return None
        return str_val.strip()
    except:
        return None

def make_jsonly_serializable(df):
    if df.empty:
        return df
    df = df.copy()
    for col in df.columns:
        try:
            if pd.api.types.is_numeric_dtype(df[col]):
                if pd.api.types.is_integer_dtype(df[col]):
                    df[col] = df[col].astype('Int64')
                else:
                    df[col] = df[col].astype(float)
            else:
                df[col] = [safe_convert_value(val) for val in df[col]]
        except Exception as e:
            st.warning(f"Error processing column '{col}': {e}")
            df[col] = [str(val) if val is not None else None for val in df[col]]
    return df.reset_index(drop=True)

def find_table_end(df, start_idx):
    for i in range(start_idx, len(df)):
        row_text = ' '.join(str(cell) for cell in df.iloc[i].values if pd.notna(cell)).upper()
        if any(term in row_text for term in ['TOTAL SALES', 'GRAND TOTAL', 'OVERALL TOTAL']):
            return i + 1  # Include the "TOTAL SALES" row
    return len(df)

# Matplotlib Chart Creation for PPT
def create_matplotlib_chart(data, x_col, y_col, chart_type, title, color_override=None):
    """Create matplotlib chart and return as BytesIO object for PPT insertion."""
    try:
        # Set style and figure size
        plt.style.use('seaborn-v0_8')
        
        # Special handling for Product Performance charts with long x-axis labels
        is_product_performance = "Product Performance" in title
        
        if is_product_performance:
            # Wider figure for product performance charts
            fig, ax = plt.subplots(figsize=(16, 10))  # Increased height for better fit
        else:
            fig, ax = plt.subplots(figsize=(14, 10))  # Increased size for better layout
            
        # Determine colors
        if color_override:
            if color_override == '#FF8C00':
                main_color = '#FF8C00'  # Orange for Act
                palette = ['#FF8C00', '#FFB347', '#FF7F00', '#FF6347', '#FF4500']
            else:
                main_color = color_override
                palette = [color_override]
        else:
            main_color = '#2E86AB'  # Default blue
            palette = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#556B2F']
        
        if chart_type.lower() == 'pie':
            # Pie chart
            colors = sns.color_palette("Set3", len(data))
            wedges, texts, autotexts = ax.pie(
                data[y_col], 
                labels=data[x_col], 
                autopct='%1.1f%%',
                colors=colors,
                startangle=90,
                textprops={'fontsize': 12, 'fontweight': 'bold'}
            )
            
            # Make percentage text more visible
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontweight('bold')
                autotext.set_fontsize(11)
                
        elif chart_type.lower() == 'line':
            # Line chart - NO VALUE LABELS
            ax.plot(data[x_col], data[y_col], 
                   marker='o', linewidth=3, markersize=8, 
                   color=main_color, markerfacecolor='white', 
                   markeredgewidth=2, markeredgecolor=main_color)
            
            # Format y-axis
            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'{x:,.0f}'))
            ax.grid(True, alpha=0.3)
            
        else:  # Bar chart (default)
            # Handle grouped data for budget vs actual
            if 'Metric' in data.columns and len(data['Metric'].unique()) > 1:
                # Grouped bar chart - NO VALUE LABELS
                metrics = data['Metric'].unique()
                x_pos = range(len(data[x_col].unique()))
                width = 0.35
                
                for i, metric in enumerate(metrics):
                    metric_data = data[data['Metric'] == metric]
                    color = '#2E86AB' if 'budget' in metric.lower() else '#FF8C00'
                    ax.bar([x + width*i for x in x_pos], 
                          metric_data[y_col], 
                          width, 
                          label=metric, 
                          color=color,
                          alpha=0.8,
                          edgecolor='darkblue',
                          linewidth=1)
                
                ax.set_xticks([x + width/2 for x in x_pos])
                ax.set_xticklabels(data[x_col].unique(), rotation=0)  # Force straight labels
                ax.legend(fontsize=12, loc='upper right')
            else:
                # Regular bar chart - NO VALUE LABELS
                bars = ax.bar(data[x_col], data[y_col], 
                             color=main_color, alpha=0.8, 
                             edgecolor='darkblue', linewidth=1)
                
                # NO VALUE LABELS ON BARS FOR PPT
            
            # Format y-axis
            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'{x:,.0f}'))
            ax.grid(True, axis='y', alpha=0.3)
        
        # Set title and labels
        ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
        ax.set_xlabel(x_col, fontsize=14, fontweight='bold')
        if chart_type.lower() != 'pie':
            ax.set_ylabel(y_col, fontsize=14, fontweight='bold')
        
        # Special handling for x-axis labels in Product Performance charts
        if is_product_performance and chart_type.lower() != 'pie':
            # Rotate labels at 45 degrees and adjust alignment
            plt.xticks(rotation=45, ha='right', fontsize=10)
            # Use manual layout adjustment instead of tight_layout for better control
            plt.subplots_adjust(left=0.1, right=0.95, top=0.9, bottom=0.3)
        else:
            # FORCE STRAIGHT X-AXIS LABELS FOR ALL OTHER PPT CHARTS
            plt.xticks(rotation=0, ha='center', fontsize=10)
            # Use manual layout adjustment for consistent spacing
            plt.subplots_adjust(left=0.1, right=0.95, top=0.9, bottom=0.15)
        
        # Save to BytesIO with manual layout instead of tight_layout
        img_buffer = BytesIO()
        plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight', 
                   facecolor='white', edgecolor='none', pad_inches=0.3)
        img_buffer.seek(0)
        
        # Clear the figure
        plt.clf()
        plt.close()
        
        return img_buffer
        
    except Exception as e:
        # Suppress the tight_layout warning and continue
        if "tight_layout" in str(e):
            # Try to save without tight layout
            try:
                img_buffer = BytesIO()
                plt.savefig(img_buffer, format='png', dpi=300, 
                           facecolor='white', edgecolor='none', pad_inches=0.3)
                img_buffer.seek(0)
                plt.clf()
                plt.close()
                return img_buffer
            except:
                pass
        
        st.warning(f"Chart generation warning (non-critical): {str(e)}")
        plt.clf()
        plt.close()
        return None

def add_data_table_to_slide(slide, data, x_col, y_col):
    """Add a data table to a PowerPoint slide."""
    try:
        # Prepare data for table
        if len(data.columns) > 10:
            # Show only first 10 columns if too many
            display_data = data.iloc[:, :10]
        else:
            display_data = data
        
        # Limit rows to prevent overcrowding
        if len(display_data) > 15:
            display_data = display_data.head(15)
        
        rows = len(display_data) + 1  # +1 for header
        cols = len(display_data.columns)
        
        # Add table
        table = slide.shapes.add_table(
            rows, cols, 
            Inches(0.5), Inches(1.5), 
            Inches(9), Inches(5)
        ).table
        
        # Set column headers
        for col_idx, col_name in enumerate(display_data.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Header background color
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(47, 117, 181)  # Blue header
            
            # Header text color
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add data rows
        for row_idx, (_, row) in enumerate(display_data.iterrows()):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                if pd.isna(value):
                    cell.text = ""
                elif isinstance(value, (int, float)):
                    cell.text = f"{value:,.2f}" if isinstance(value, float) else f"{value:,}"
                else:
                    cell.text = str(value)
                
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Alternating row colors
                if row_idx % 2 == 0:
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(242, 242, 242)  # Light gray
        
        # Add note if data was truncated
        if len(data) > 15 or len(data.columns) > 10:
            note_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.8), Inches(9), Inches(0.5)
            )
            note_frame = note_shape.text_frame
            note_frame.text = f"Note: Showing first {min(15, len(data))} rows and {min(10, len(data.columns))} columns of {len(data)} total rows."
            note_frame.paragraphs[0].font.size = Pt(10)
            note_frame.paragraphs[0].font.italic = True
            
    except Exception as e:
        st.error(f"Could not add data table to slide: {str(e)}")

def create_plotly_chart(data, x_col, y_col, chart_type, title, color_override=None):
    default_color = color_override if color_override else '#2E86AB'
    
    layout_config = {
        'title': {
            'text': title,
            'x': 0.5,
            'xanchor': 'center',
            'font': {'size': 20, 'family': 'Arial, sans-serif'}
        },
        'font': {'size': 14, 'family': 'Arial, sans-serif'},
        'plot_bgcolor': 'white',
        'paper_bgcolor': 'white',
        'height': 600,
        'margin': {'l': 80, 'r': 80, 't': 100, 'b': 80},
        'hovermode': 'closest',  # Show closest data point on hover
        'showlegend': True if chart_type == 'pie' else False
    }
    
    if chart_type == 'bar':
        fig = px.bar(data, x=x_col, y=y_col, color_discrete_sequence=[default_color])
        fig.update_traces(
            hovertemplate='<b>%{x}</b><br>Value: %{y:,.0f}<extra></extra>'
        )
        fig.update_xaxes(
            title_font={'size': 16, 'family': 'Arial, sans-serif'},
            tickfont={'size': 14},
            tickangle=45 if max(len(str(x)) for x in data[x_col]) > 10 else 0
        )
        fig.update_yaxes(
            title_font={'size': 16, 'family': 'Arial, sans-serif'},
            tickfont={'size': 14}
        )
        
    elif chart_type == 'line':
        fig = px.line(data, x=x_col, y=y_col, markers=True, color_discrete_sequence=[default_color])
        fig.update_traces(
            line={'width': 4}, 
            marker={'size': 10},
            hovertemplate='<b>%{x}</b><br>Value: %{y:,.0f}<extra></extra>'
        )
        fig.update_xaxes(
            title_font={'size': 16, 'family': 'Arial, sans-serif'},
            tickfont={'size': 14},
            tickangle=45 if max(len(str(x)) for x in data[x_col]) > 10 else 0
        )
        fig.update_yaxes(
            title_font={'size': 16, 'family': 'Arial, sans-serif'},
            tickfont={'size': 14}
        )
        
    elif chart_type == 'pie':
        fig = px.pie(data, values=y_col, names=x_col)
        fig.update_traces(
            textposition='inside',
            textinfo='percent',
            hovertemplate='<b>%{label}</b><br>Value: %{value:,.0f}<br>Percentage: %{percent:.1%}<extra></extra>',
            textfont={'size': 14, 'family': 'Arial, sans-serif'}
        )
    
    fig.update_layout(**layout_config)
    
    config = {
        'displayModeBar': True,
        'displaylogo': False,
        'modeBarButtonsToRemove': ['pan2d', 'lasso2d', 'select2d'],
        'toImageButtonOptions': {
            'format': 'png',
            'filename': 'chart',
            'height': 800,
            'width': 1200,
            'scale': 2
        }
    }
    
    return fig, config

def create_ppt_with_chart(title, chart_data, x_col, y_col, chart_type='bar', color_override=None, selected_filter=None):
    """Create PowerPoint presentation with matplotlib charts (Streamlit Cloud compatible)."""
    try:
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        
        # Add selected filter to title if provided
        if selected_filter and selected_filter != "Select All":
            title = f"{title} - {selected_filter}"
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = title
        
        if chart_data is None or chart_data.empty:
            st.error(f"Error: No data provided for {title}.")
            error_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
            ef = error_box.text_frame
            ef.text = "Error: No data available"
            ppt_bytes = BytesIO()
            ppt.save(ppt_bytes)
            ppt_bytes.seek(0)
            return ppt_bytes
        
        if y_col not in chart_data.columns:
            st.error(f"Error: Column {y_col} not found in data for {title}.")
            error_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
            ef = error_box.text_frame
            ef.text = f"Error: Column {y_col} not found"
            ppt_bytes = BytesIO()
            ppt.save(ppt_bytes)
            ppt_bytes.seek(0)
            return ppt_bytes
        
        if not pd.api.types.is_numeric_dtype(chart_data[y_col]):
            st.error(f"Error: Column {y_col} is not numeric for {title}. Cannot create chart.")
            error_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
            ef = error_box.text_frame
            ef.text = f"Error: No numeric data available for {y_col}"
            ppt_bytes = BytesIO()
            ppt.save(ppt_bytes)
            ppt_bytes.seek(0)
            return ppt_bytes
        
        if chart_type == 'pie':
            chart_data = chart_data[chart_data[y_col] > 0].copy()
            if chart_data.empty:
                st.warning(f"No positive values available for pie chart in {title}. Skipping chart creation.")
                error_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
                ef = error_box.text_frame
                ef.text = f"No positive values available for {y_col} pie chart"
                ppt_bytes = BytesIO()
                ppt.save(ppt_bytes)
                ppt_bytes.seek(0)
                return ppt_bytes
        
        # Try matplotlib first (for cloud compatibility)
        try:
            img_buffer = create_matplotlib_chart(chart_data, x_col, y_col, chart_type, title, color_override)
            if img_buffer:
                slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
            else:
                raise Exception("Matplotlib chart generation failed")
        except Exception as matplotlib_error:
            # Fallback to Plotly if matplotlib fails
            try:
                fig, _ = create_plotly_chart(chart_data, x_col, y_col, chart_type, title, color_override)
                img_bytes = fig.to_image(format="png", width=1200, height=800, scale=2)
                img_buffer = BytesIO(img_bytes)
                slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
            except Exception as plotly_error:
                # Final fallback with text placeholder
                st.warning(f"Could not add chart image to PPT: {plotly_error}")
                text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                text_frame = text_box.text_frame
                text_frame.text = f"Chart: {title}\n(Image generation not available in this environment)"
        
        ppt_bytes = BytesIO()
        ppt.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes
        
    except Exception as e:
        st.error(f"Could not create PPT: {str(e)}")
        return None

def create_master_ppt_with_matplotlib(all_data, table_name, selected_sheet, visual_type, selected_filter=None):
    """Create master PPT with all visualizations using matplotlib - CHARTS ONLY, NO TABLES."""
    try:
        master_ppt = Presentation()
        
        # Add title slide
        title_slide_layout = master_ppt.slide_layouts[0]
        title_slide = master_ppt.slides.add_slide(title_slide_layout)
        
        # Add selected filter to title if provided
        if selected_filter and selected_filter != "Select All":
            title_slide.shapes.title.text = f"Complete Analysis Report - {table_name} - {selected_filter}"
        else:
            title_slide.shapes.title.text = f"Complete Analysis Report - {table_name}"
            
        if title_slide.shapes.placeholders[1]:
            title_slide.shapes.placeholders[1].text = f"Sheet: {selected_sheet}\nGenerated on: {datetime.now().strftime('%Y-%m-%d')}"
        
        for label, data in all_data:
            if data is not None and (not isinstance(data, pd.DataFrame) or not data.empty):
                try:
                    # Determine chart configuration based on label
                    if label == "Budget vs Actual":
                        if 'Month' in data.columns and 'Value' in data.columns and 'Metric' in data.columns:
                            chart_data = data
                            x_col = "Month"
                            y_col = "Value"
                            chart_type = visual_type.lower().replace(" chart", "")
                            color_override = None  # Will use default colors for grouped data
                        else:
                            continue
                    elif label == "Branch Performance":
                        if 'Branch' in data.columns and 'Performance' in data.columns:
                            chart_data = data
                            x_col = "Branch"
                            y_col = "Performance"
                            chart_type = visual_type.lower().replace(" chart", "")
                            color_override = None
                        elif len(data.columns) >= 2:
                            x_col = data.columns[0]
                            y_col = data.columns[1]
                            chart_data = data
                            chart_type = visual_type.lower().replace(" chart", "")
                            color_override = None
                        else:
                            continue
                    elif label == "Product Performance":
                        if 'Product' in data.columns and 'Performance' in data.columns:
                            chart_data = data
                            x_col = "Product"
                            y_col = "Performance"
                            chart_type = visual_type.lower().replace(" chart", "")
                            color_override = None
                        elif len(data.columns) >= 2:
                            x_col = data.columns[0]
                            y_col = data.columns[1]
                            chart_data = data
                            chart_type = visual_type.lower().replace(" chart", "")
                            color_override = None
                        else:
                            continue
                    elif label in ["Branch Monthwise", "Product Monthwise"]:
                        if 'Month' in data.columns and 'Value' in data.columns:
                            chart_data = data
                            x_col = "Month"
                            y_col = "Value"
                            chart_type = visual_type.lower().replace(" chart", "")
                            color_override = None
                        else:
                            continue
                    elif "YTD" in label:
                        if 'Period' in data.columns:
                            chart_data = data
                            x_col = "Period"
                            y_col = label.replace("YTD ", "")
                            chart_type = visual_type.lower().replace(" chart", "")
                            color_override = '#FF8C00' if 'Act' in label else None
                        elif len(data.columns) >= 2:
                            chart_data = data
                            x_col = data.columns[0]
                            y_col = data.columns[1]
                            chart_type = visual_type.lower().replace(" chart", "")
                            color_override = '#FF8C00' if 'Act' in label else None
                        else:
                            continue
                    else:
                        if "Month" in data.columns:
                            label_clean = label.replace(",", "").replace(" ", "")
                            if label_clean in data.columns:
                                chart_data = data
                                x_col = "Month"
                                y_col = label_clean
                                chart_type = visual_type.lower().replace(" chart", "")
                                color_override = '#FF8C00' if label == "Act" else None
                            elif "Value" in data.columns:
                                chart_data = data
                                x_col = "Month"
                                y_col = "Value"
                                chart_type = visual_type.lower().replace(" chart", "")
                                color_override = '#FF8C00' if label == "Act" else None
                        elif "Period" in data.columns:
                            label_clean = label.replace(",", "").replace(" ", "")
                            if label_clean in data.columns:
                                chart_data = data
                                x_col = "Period"
                                y_col = label_clean
                                chart_type = visual_type.lower().replace(" chart", "")
                                color_override = '#FF8C00' if label == "Act" else None
                            elif "Value" in data.columns:
                                chart_data = data
                                x_col = "Period"
                                y_col = "Value"
                                chart_type = visual_type.lower().replace(" chart", "")
                                color_override = '#FF8C00' if label == "Act" else None
                        else:
                            continue
                    
                    if chart_data is None or x_col is None or y_col is None:
                        continue
                    
                    # Add ONLY chart slide - NO DATA TABLES
                    chart_slide_layout = master_ppt.slide_layouts[6]  # Blank layout
                    slide = master_ppt.slides.add_slide(chart_slide_layout)
                    
                    # Add slide title with filter if applicable
                    if selected_filter and selected_filter != "Select All":
                        slide_title = f"{label} Analysis - {table_name} - {selected_filter}"
                    else:
                        slide_title = f"{label} Analysis - {table_name}"
                    
                    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                    title_frame = title_shape.text_frame
                    title_frame.text = slide_title
                    title_frame.paragraphs[0].font.size = Pt(24)
                    title_frame.paragraphs[0].font.bold = True
                    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # Generate matplotlib chart with straight x-axis labels
                    img_buffer = create_matplotlib_chart(chart_data, x_col, y_col, chart_type, 
                                                        slide_title, color_override)
                    
                    if img_buffer:
                        # Add chart image to slide - LARGER SIZE since no table needed
                        slide.shapes.add_picture(
                            img_buffer, 
                            Inches(0.5), 
                            Inches(1.2), 
                            width=Inches(9), 
                            height=Inches(6.5)  # Larger chart without table
                        )
                    else:
                        # Add text placeholder if image generation fails
                        text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
                        text_frame = text_box.text_frame
                        text_frame.text = f"Chart: {label}\n(Image generation not available)"
                        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                except Exception as e:
                    st.warning(f"Error creating chart for {label}: {e}")
                    continue
        
        # Save to BytesIO
        master_ppt_buffer = BytesIO()
        master_ppt.save(master_ppt_buffer)
        master_ppt_buffer.seek(0)
        
        return master_ppt_buffer.getvalue()
        
    except Exception as e:
        st.error(f"Could not create master PPT: {str(e)}")
        return None

def ensure_numeric_data(data, y_col):
    if y_col not in data.columns:
        return False
    try:
        data[y_col] = pd.to_numeric(data[y_col].astype(str).str.replace(',', ''), errors='coerce')
        data.dropna(subset=[y_col], inplace=True)
    except Exception as e:
        st.warning(f"Failed to convert {y_col} to numeric: {e}")
        return False
    return not data.empty

def display_visualization(tab, label, data, x_col, y_col, visual_type, color_override=None):
    with tab:
        if data is None or data.empty:
            st.warning(f"No data available for {label}")
            return
        
        if not ensure_numeric_data(data, y_col):
            st.warning(f"No numeric data available for {label}")
            return None
        
        if visual_type == "Pie Chart":
            data = data[data[y_col] > 0]
            if data.empty:
                st.warning(f"No positive values available for {label} pie chart")
                return None
        
        st.markdown(f"### {label} - {table_name}")
        
        chart_type_map = {
            "Bar Chart": "bar",
            "Line Chart": "line", 
            "Pie Chart": "pie"
        }
        
        fig, config = create_plotly_chart(
            data, x_col, y_col, 
            chart_type_map[visual_type], 
            f"{label} - {table_name}",
            color_override
        )
        
        st.plotly_chart(fig, use_container_width=True, config=config)
        
        with st.expander("📊 View Data Table"):
            st.dataframe(data, use_container_width=True)
        
        optimize_memory()

def extract_month_year(col_name):
    """Extract clean month-year format from column names."""
    col_str = str(col_name).strip()
    
    # Remove common prefixes
    col_str = re.sub(r'^(Gr[-\s]*|Ach[-\s]*|Act[-\s]*|Budget[-\s]*|LY[-\s]*)', '', col_str, flags=re.IGNORECASE)
    
    # Extract month-year pattern
    month_year_match = re.search(r'(\w{3,})[-–\s]*(\d{2})', col_str, re.IGNORECASE)
    if month_year_match:
        month, year = month_year_match.groups()
        return f"{month.capitalize()}-{year}"
    
    return col_str

# File uploader
uploaded_file = st.sidebar.file_uploader("📂 Upload Excel File", type=["xlsx"])

if uploaded_file:
    try:
        optimize_memory()
        
        file_size = uploaded_file.size
        chunk_size = 1024 * 1024
        
        if file_size > 10 * 1024 * 1024:
            st.warning("Large file detected. Processing in chunks for better performance...")
            # Use BytesIO instead of temporary file for cloud compatibility
            file_buffer = BytesIO(uploaded_file.getbuffer())
            xls = pd.ExcelFile(file_buffer)
        else:
            xls = pd.ExcelFile(uploaded_file)
            
        sheet_names = xls.sheet_names
        selected_sheet = st.sidebar.selectbox("📄 Select a Sheet", sheet_names)
        
        df_sheet = pd.read_excel(xls, sheet_name=selected_sheet, header=None, nrows=1000)
        
        if df_sheet.shape[1] < 10 and df_sheet.iloc[:, 0].astype(str).str.len().max() > 200:
            try:
                df_sheet_alt = pd.read_excel(xls, sheet_name=selected_sheet, header=None, engine='openpyxl', nrows=1000)
                if df_sheet_alt.shape[1] > df_sheet.shape[1]:
                    df_sheet = df_sheet_alt
                    st.info("✅ Improved data structure using alternative reading method")
            except:
                pass
                
        if df_sheet.shape[1] < 20:
            new_data = []
            months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
            metrics = ['Budget', 'LY', 'Act', 'Gr', 'Ach', 'YTD']
            year_pattern = r'\d{2,4}(?:[-–]\d{2,4})?'
            
            for idx, row in df_sheet.iterrows():
                if pd.notna(row.iloc[0]):
                    row_text = str(row.iloc[0]).strip()
                    if any(metric in row_text for metric in metrics) or re.search(r'SALES\s*(in\s*(MT|Value|Ton[n]?age))?', row_text, re.IGNORECASE):
                        patterns = []
                        patterns.append(r'SALES\s*in\s*(MT|Value|Ton[n]?age)', re.IGNORECASE)
                        for metric in metrics:
                            for month in months:
                                patterns.append(rf'{metric}[-–\s]*{month}[-–\s]*{year_pattern}', re.IGNORECASE)
                            patterns.append(rf'{metric}[-–\s]*YTD[-–\s]*{year_pattern}\s*\([^)]*\)', re.IGNORECASE)
                            patterns.append(rf'YTD[-–\s]*{year_pattern}\s*\([^)]*\)\s*{metric}', re.IGNORECASE)
                        
                        positions = []
                        for pattern in patterns:
                            for match in re.finditer(pattern, row_text):
                                positions.append((match.start(), match.group()))
                        positions.sort()
                        parts = [item[1].strip() for item in positions]
                        
                        if len(parts) < 5:
                            parts = [part.strip() for part in row_text.split() if part.strip()]
                        
                        new_data.append(parts)
                    else:
                        new_data.append(row_text.split())
                else:
                    new_data.append([])
            
            if new_data:
                max_cols = max(len(row) for row in new_data)
                for row in new_data:
                    while len(row) < max_cols:
                        row.append(None)
                df_sheet = pd.DataFrame(new_data)
        
    except Exception as e:
        st.error(f"Error reading Excel file: {e}")
        st.stop()

    sheet_index = sheet_names.index(selected_sheet)
    is_first_sheet = sheet_index == 0
    is_sales_monthwise = 'sales analysis month wise' in selected_sheet.lower() or ('sales' in selected_sheet.lower() and 'month' in selected_sheet.lower())

    table_name = ""

    if is_first_sheet:
        st.subheader("📋 First Sheet - Table Detection")
        
        table1_start = None
        table2_start = None
        for i in range(len(df_sheet)):
            row_text = ' '.join(str(cell) for cell in df_sheet.iloc[i].values if pd.notna(cell))
            if re.search(r'\bsales\s*in\s*mt\b', row_text, re.IGNORECASE) and table1_start is None:
                table1_start = i
            elif re.search(r'\bsales\s*in\s*(value|tonnage|tonage)\b', row_text, re.IGNORECASE) and table1_start is not None and table2_start is None:
                table2_start = i
        
        table_options = []
        if table1_start is not None:
            table_options.append("Table 1: SALES IN MT")
        if table2_start is not None:
            table_options.append("Table 2: SALES IN VALUE")
        
        if table_options:
            table_choice = st.sidebar.radio("📌 Select Table", table_options, key="first_sheet_table_select")
            table_name = table_choice
            
            if table_choice == "Table 1: SALES IN MT" and table1_start is not None:
                st.write("### Table 1: SALES IN MT")
                table1_end = table2_start if table2_start is not None else len(df_sheet)
                table1 = df_sheet.iloc[table1_start:table1_end].dropna(how='all').reset_index(drop=True)

                if not table1.empty:
                    header_row_idx = None
                    for i in range(min(5, len(table1))):
                        row_text = ' '.join(str(cell) for cell in table1.iloc[i].values if pd.notna(cell))
                        if re.search(r'\b(?:budget|ly|act|gr|ach)\b', row_text, re.IGNORECASE):
                            header_row_idx = i
                            break
                    
                    if header_row_idx is not None:
                        header_row = table1.iloc[header_row_idx]
                        new_columns = [str(val).strip() if pd.notna(val) else f'Unnamed_{i}' 
                                      for i, val in enumerate(header_row)]
                        table1.columns = new_columns
                        table1 = table1.iloc[header_row_idx + 1:].reset_index(drop=True)
                        
                        if 2 <= sheet_index <= 4:
                            if not table1.empty:
                                table1 = table1.drop(index=0).reset_index(drop=True)
                            else:
                                st.warning("Table 1 is empty after processing, cannot delete first row.")
                        
                        if sheet_index == 1 and not table1.empty:
                            table1 = table1[~table1[table1.columns[0]].str.contains('REGIONS', case=False, na=False, regex=True)].reset_index(drop=True)
                        
                        if not table1.empty:
                            table1 = make_jsonly_serializable(table1)
                            st.dataframe(table1, use_container_width=True)
                            csv1 = table1.to_csv(index=False).encode('utf-8')
                            st.download_button(
                                "⬇️ Download Table 1 as CSV", 
                                csv1, 
                                "sales_in_mt.csv", 
                                "text/csv",
                                key="download_table1_csv"
                            )
                        else:
                            st.warning("No data available for Table 1 after processing.")
                    else:
                        st.error("Could not find column headers for Table 1.")
                        st.dataframe(table1)
            
            elif table_choice == "Table 2: SALES IN VALUE" and table2_start is not None:
                if sheet_index >= 1 and sheet_index <= 4:
                    table2_end = find_table_end(df_sheet, table2_start)
                    table2 = df_sheet.iloc[table2_start:table2_end].dropna(how='all').reset_index(drop=True)
                else:
                    table2 = df_sheet.iloc[table2_start:].dropna(how='all').reset_index(drop=True)
                
                if not table2.empty:
                    header_row_idx = None
                    for i in range(min(5, len(table2))):
                        row_text = ' '.join(str(cell) for cell in table2.iloc[i].values if pd.notna(cell))
                        if re.search(r'\b(?:budget|ly|act|gr|ach)\b', row_text, re.IGNORECASE):
                            header_row_idx = i
                            break
                    
                    if header_row_idx is not None:
                        header_row = table2.iloc[header_row_idx]
                        new_columns = [str(val).strip() if pd.notna(val) else f'Unnamed_{i}' 
                                      for i, val in enumerate(header_row)]
                        table2.columns = new_columns
                        table2 = table2.iloc[header_row_idx + 1:].reset_index(drop=True)
                        
                        if 2 <= sheet_index <= 4:
                            if not table2.empty:
                                table2 = table2.drop(index=0).reset_index(drop=True)
                            else:
                                st.warning("Table 2 is empty after processing, cannot delete first row.")
                        
                        if sheet_index == 1 and not table2.empty:
                            table2 = table2[~table2[table2.columns[0]].str.contains('REGIONS', case=False, na=False, regex=True)].reset_index(drop=True)
                        
                        if not table2.empty:
                            table2 = make_jsonly_serializable(table2)
                            st.dataframe(table2, use_container_width=True)
                            csv2 = table2.to_csv(index=False).encode('utf-8')
                            st.download_button(
                                "⬇️ Download Table 2 as CSV", 
                                csv2, 
                                "sales_in_value.csv", 
                                "text/csv",
                                key="download_table2_csv"
                            )
                        else:
                            st.warning("No data available for Table 2 after processing.")
                    else:
                        st.error("Could not find column headers for Table 2.")
                        st.dataframe(table2)
                else:
                    st.warning("Table 2 is empty or contains no valid data.")
        else:
            st.warning("No tables ('SALES IN MT' or 'SALES IN VALUE') found in the sheet.")
            df_sheet_clean = make_jsonly_serializable(df_sheet)
            st.dataframe(df_sheet_clean, use_container_width=True)
            csv = df_sheet_clean.to_csv(index=False).encode('utf-8')
            st.download_button(
                "⬇️ Download Raw Data as CSV", 
                csv, 
                "raw_data.csv", 
                "text/csv",
                key="download_raw_data_csv"
            )
    
    else:
        is_product_analysis = ('product' in selected_sheet.lower() or 
                             'ts-pw' in selected_sheet.lower() or 
                             'ero-pw' in selected_sheet.lower())
        is_branch_analysis = 'region wise analysis' in selected_sheet.lower()

        if is_branch_analysis:
            table1_header = "Sales in MT"
            table2_header = "Sales in Value"
        else:
            table1_header = "Sales in Tonage"
            table2_header = "Sales in Value"

        def extract_tables(df, table1_header, table2_header):
            table1_idx, table2_idx = None, None
            for i in range(len(df)):
                row_text = ' '.join(df.iloc[i].astype(str).str.lower().tolist())
                if table1_idx is None and table1_header.lower() in row_text:
                    table1_idx = i
                elif table2_idx is None and table2_header.lower() in row_text and i > (table1_idx or 0):
                    table2_idx = i
            return table1_idx, table2_idx

        idx1, idx2 = extract_tables(df_sheet, table1_header, table2_header)

        if idx1 is None:
            st.error(f"❌ Could not locate '{table1_header}' header in the sheet.")
        else:
            if sheet_index >= 1 and sheet_index <= 4:
                table1_end = find_table_end(df_sheet, idx1 + 1)
                if idx2 is not None:
                    table1_end = min(table1_end, idx2)  # Ensure table1_end includes "TOTAL SALES"
            else:
                table1_end = idx2 if idx2 is not None else len(df_sheet)
            
            # Extract Table 1, including the header row
            table1 = df_sheet.iloc[idx1:table1_end].dropna(how='all').reset_index(drop=True)
            
            if not table1.empty:
                # Use the first row as the header
                table1.columns = table1.iloc[0].apply(lambda x: str(x).strip() if pd.notna(x) else '')
                table1 = table1.iloc[1:].reset_index(drop=True)

            if idx2 is not None:
                if sheet_index >= 1 and sheet_index <= 4:
                    table2_end = find_table_end(df_sheet, idx2 + 1)
                else:
                    table2_end = len(df_sheet)
                
                # Extract Table 2, including the header row
                table2 = df_sheet.iloc[idx2:table2_end].dropna(how='all').reset_index(drop=True)
                
                if not table2.empty:
                    # Use the first row as the header
                    table2.columns = table2.iloc[0].apply(lambda x: str(x).strip() if pd.notna(x) else '')
                    table2 = table2.iloc[1:].reset_index(drop=True)
            else:
                table2 = None

            table_options = [f"Table 1: {table1_header.upper()}"]
            if table2 is not None:
                table_options.append(f"Table 2: {table2_header.upper()}")
            table_choice = st.sidebar.radio("📌 Select Table", table_options)
            table_name = table_choice
            table_df = table1 if table_choice == table_options[0] else table2

            table_df = make_jsonly_serializable(table_df)
            table_df.columns = table_df.columns.map(str)

            def rename_columns(columns):
                renamed = []
                ytd_base = None
                prev_month = None
                prev_year = None
                for col in columns:
                    col_clean = col.strip().replace(",", "").replace("–", "-")
                    ytd_act_match = re.search(r'(YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))\s*Act', col_clean, re.IGNORECASE)
                    if ytd_act_match:
                        ytd_base = ytd_act_match.group(1).replace("–", "-")
                        renamed.append(f"Act-{ytd_base}")
                        continue
                    ytd_gr_match = re.search(r'(Gr-YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))', col_clean, re.IGNORECASE)
                    if ytd_gr_match:
                        gr_ytd = ytd_gr_match.group(1).replace("–", "-")
                        renamed.append(f"Gr-{gr_ytd.split('Gr-')[1]}")
                        continue
                    ytd_ach_match = re.search(r'(Ach-YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))', col_clean, re.IGNORECASE)
                    ytd_ach_alt_match = re.search(r'(YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))\s*Ach', col_clean, re.IGNORECASE)
                    if ytd_ach_match:
                        ach_ytd = ytd_ach_match.group(1).replace("–", "-")
                        renamed.append(ach_ytd)
                        continue
                    elif ytd_ach_alt_match:
                        ytd_part = ytd_ach_alt_match.group(1).replace("–", "-")
                        renamed.append(f"Ach-{ytd_part}")
                        continue
                    monthly_match = re.search(r'(\b\w{3,})[\s-]*(\d{2})', col_clean)
                    if monthly_match:
                        prev_month, prev_year = monthly_match.groups()
                        prev_month = prev_month.capitalize()
                    if col_clean.lower().startswith("gr") and prev_month and prev_year:
                        renamed.append(f"Gr - {prev_month}-{prev_year}")
                    elif col_clean.lower().startswith("ach") and prev_month and prev_year:
                        renamed.append(f"Ach - {prev_month}-{prev_year}")
                    else:
                        renamed.append(col)
                return renamed

            table_df.columns = rename_columns(table_df.columns)

            if table_df.columns.duplicated().any():
                table_df = table_df.loc[:, ~table_df.columns.duplicated()]
                st.warning("⚠️ Duplicate columns detected and removed.")

            branch_list = []
            product_list = []

            def extract_unique_values(df, first_col, exclude_terms=None):
                valid_rows = df[df[first_col].notna()]
                unique_values = valid_rows[first_col].astype(str).str.strip().unique()
                return sorted(unique_values)

            if is_branch_analysis:
                branch_list = extract_unique_values(table_df, table_df.columns[0])
            elif is_product_analysis:
                product_list = extract_unique_values(table_df, table_df.columns[0])

            months = sorted(set(re.findall(r'\b(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\b', 
                       ' '.join(table_df.columns), flags=re.IGNORECASE)))
            years = sorted(set(re.findall(r'[-–](\d{2})\b', ' '.join(table_df.columns))))

            months_options = ["Select All"] + months
            years_options = ["Select All"] + years
            branches_options = ["Select All"] + branch_list if is_branch_analysis else []
            products_options = ["Select All"] + product_list if is_product_analysis else []

            selected_month = st.sidebar.selectbox("📅 Filter by Month", months_options, index=0)
            selected_year = st.sidebar.selectbox("📆 Filter by Year", years_options, index=0)
            selected_branch = st.sidebar.selectbox("🌍 Filter by Branch", branches_options, index=0) if is_branch_analysis else None
            selected_product = st.sidebar.selectbox("📦 Filter by Product", products_options, index=0) if is_product_analysis else None

            selected_months = months if selected_month == "Select All" else [selected_month] if selected_month else []
            selected_years = years if selected_year == "Select All" else [selected_year] if selected_year else []
            selected_branches = branch_list if selected_branch == "Select All" else [selected_branch] if selected_branch else []
            selected_products = product_list if selected_product == "Select All" else [selected_product] if selected_product else []

            filtered_df = table_df.copy()
            first_col = filtered_df.columns[0]

            if selected_branches and is_branch_analysis:
                filtered_df = filtered_df[filtered_df[first_col].astype(str).isin(selected_branches)]
            if selected_products and is_product_analysis:
                filtered_df = filtered_df[filtered_df[first_col].astype(str).isin(selected_products)]

            def column_filter(col):
                col_str = str(col).lower().replace(",", "").replace("–", "-")
                if "ytd" in col_str:
                    return any(f"-{y}" in col_str for y in selected_years) if selected_years else True
                month_match = any(m.lower() in col_str for m in selected_months)
                year_match = any(f"-{y}" in col_str for y in selected_years) if selected_years else True
                return month_match and year_match

            visual_cols = [col for col in table_df.columns if column_filter(col)]
            display_df = filtered_df[[first_col] + visual_cols] if visual_cols else filtered_df[[first_col]]

            display_df = make_jsonly_serializable(display_df)

            def convert_to_numeric(series):
                try:
                    return pd.to_numeric(series.astype(str).str.replace(',', ''), errors='coerce')
                except:
                    return series

            formatted_df = display_df.copy()
            numeric_cols = []
            for col in formatted_df.columns:
                if col == formatted_df.columns[0]:
                    continue
                formatted_df[col] = convert_to_numeric(formatted_df[col])
                if pd.api.types.is_numeric_dtype(formatted_df[col]):
                    numeric_cols.append(col)
                    formatted_df[col] = formatted_df[col].round(2)

            style_dict = {col: "{:.2f}" for col in numeric_cols}
            st.subheader("📋 Filtered Table View")
            st.dataframe(formatted_df.style.format(style_dict, na_rep="-"), use_container_width=True)

            csv = display_df.to_csv(index=False).encode('utf-8')
            st.download_button(
                "⬇️ Download Filtered Data as CSV", 
                csv, 
                "filtered_data.csv",
                "text/csv",
                key="download_filtered_data_csv"
            )

            st.sidebar.markdown("---")
            st.sidebar.subheader("📊 Visualization Options")
            
            visual_type = st.sidebar.selectbox(
                "Select Visualization Type",
                ["Bar Chart", "Pie Chart", "Line Chart"],
                index=0,
                key="visualization_type_select"
            )

            tabs = st.tabs([
                "📊 Budget vs Actual", "📊 Budget", "📊 LY", "📊 Act", "📊 Gr", "📊 Ach", 
                "📈 YTD Budget", "📈 YTD LY", "📈 YTD Act", "📈 YTD Gr", "📈 YTD Ach", 
                "🌍 Branch Performance", "🌍 Branch Monthwise", 
                "📦 Product Performance", "📦 Product Monthwise"
            ])
            tab_names = [
                "Budget vs Actual", "Budget", "LY", "Act", "Gr", "Ach",
                "YTD Budget", "YTD LY", "YTD Act", "YTD Gr", "YTD Ach",
                "Branch Performance", "Branch Monthwise",
                "Product Performance", "Product Monthwise"
            ]
            tabs_dict = dict(zip(tab_names, tabs))

            def plot_budget_vs_actual(tab, visual_type):
                with tab:
                    budget_cols = [col for col in table_df.columns 
                                  if str(col).lower().startswith('budget') and 'ytd' not in str(col).lower()
                                  and column_filter(col)]
                    act_cols = [col for col in table_df.columns 
                                if str(col).lower().startswith('act') and 'ytd' not in str(col).lower()
                                and column_filter(col)]

                    if not (budget_cols and act_cols):
                        st.info("No matching Budget or Act columns found for comparison")
                        return None

                    budget_months = [re.search(r'(\w{3,})[-–](\d{2})', str(col), re.IGNORECASE) for col in budget_cols]
                    act_months = [re.search(r'(\w{3,})[-–](\d{2})', str(col), re.IGNORECASE) for col in act_cols]
                    common_months = set((m.group(1), m.group(2)) for m in budget_months if m) & \
                                    set((m.group(1), m.group(2)) for m in act_months if m)

                    if not common_months:
                        st.info("No common months found for Budget vs Actual comparison")
                        return None

                    selected_budget_cols = []
                    selected_act_cols = []
                    for month, year in common_months:
                        for col in budget_cols:
                            if re.search(rf'\b{month}[-–]{year}\b', str(col), re.IGNORECASE):
                                selected_budget_cols.append(col)
                        for col in act_cols:
                            if re.search(rf'\b{month}[-–]{year}\b', str(col), re.IGNORECASE):
                                selected_act_cols.append(col)

                    chart_data = filtered_df[[first_col] + selected_budget_cols + selected_act_cols].copy()

                    for col in selected_budget_cols + selected_act_cols:
                        chart_data[col] = pd.to_numeric(chart_data[col].astype(str).str.replace(',', ''), 
                                                       errors='coerce')

                    chart_data = chart_data.dropna()

                    if chart_data.empty:
                        st.warning("No valid numeric data available for Budget vs Act comparison")
                        return None

                    if visual_type == "Pie Chart":
                        budget_total = chart_data[selected_budget_cols].sum().sum()
                        act_total = chart_data[selected_act_cols].sum().sum()
                        pie_data = pd.DataFrame({
                            "Metric": ["Budget", "Act"],
                            "Value": [budget_total, act_total]
                        })
                        pie_data = pie_data[pie_data["Value"] > 0]
                        if pie_data.empty:
                            st.warning("No valid data for Budget vs Actual pie chart after aggregation")
                            return None
                        display_visualization(tab, "Budget vs Actual", pie_data, "Metric", "Value", visual_type)
                        ppt_type = 'pie'
                        chart_data_for_ppt = pie_data
                        x_col_for_ppt = "Metric"
                    else:
                        chart_data_melt = chart_data.melt(id_vars=first_col, 
                                                         var_name="Month_Metric", 
                                                         value_name="Value")
                        chart_data_melt['Metric'] = chart_data_melt['Month_Metric'].apply(
                            lambda x: 'Budget' if 'budget' in str(x).lower() else 'Act'
                        )
                        chart_data_melt['Month'] = chart_data_melt['Month_Metric'].apply(
                            lambda x: re.search(r'(\w{3,})[-–](\d{2})', str(x), re.IGNORECASE).group(0) 
                                      if re.search(r'(\w{3,})[-–](\d{2})', str(x), re.IGNORECASE) else x
                        )
                        
                        chart_data_melt = make_jsonly_serializable(chart_data_melt)
                        chart_data_agg = chart_data_melt.groupby(['Month', 'Metric'])['Value'].sum().reset_index()
                        
                        if chart_data_agg.empty or 'Value' not in chart_data_agg.columns:
                            st.warning("Aggregation failed: No valid data for Budget vs Actual comparison")
                            return None
                        
                        chart_data_agg['Value'] = pd.to_numeric(chart_data_agg['Value'], errors='coerce')
                        if chart_data_agg['Value'].isna().all():
                            st.warning("No numeric values available in aggregated data for Budget vs Actual")
                            return None
                        
                        if not ensure_numeric_data(chart_data_agg, 'Value'):
                            st.warning("No numeric data available for Budget vs Actual comparison")
                            return None
                        
                        # Define fiscal year month order
                        fiscal_month_order = ['Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar']
                        
                        # Extract month from Month string (e.g., "Apr-23" -> "Apr")
                        chart_data_agg['MonthOnly'] = chart_data_agg['Month'].str.extract(r'^(\w{3})', expand=False)
                        
                        # Create a custom sort key based on fiscal year order
                        def get_fiscal_sort_key(row):
                            month = row['MonthOnly']
                            year = re.search(r'[-–](\d{2})', row['Month']).group(1) if re.search(r'[-–](\d{2})', row['Month']) else '00'
                            try:
                                month_idx = fiscal_month_order.index(month)
                                # For months Apr-Dec, use current year, for Jan-Mar use next year
                                fiscal_year = int(year) if month_idx < 9 else int(year) + 1
                                return (fiscal_year, month_idx)
                            except ValueError:
                                return (0, 99)
                        
                        # Sort the data by fiscal year
                        chart_data_agg['SortKey'] = chart_data_agg.apply(get_fiscal_sort_key, axis=1)
                        chart_data_agg = chart_data_agg.sort_values('SortKey')
                        
                        st.markdown(f"### Budget vs Actual Comparison - {table_name}")
                        
                        fig = go.Figure()
                        budget_data = chart_data_agg[chart_data_agg['Metric'] == 'Budget']
                        act_data = chart_data_agg[chart_data_agg['Metric'] == 'Act']
                        
                        if not budget_data.empty:
                            fig.add_trace(go.Bar(
                                x=budget_data['Month'],
                                y=budget_data['Value'],
                                name='Budget',
                                marker_color='#2E86AB',
                                hovertemplate='<b>%{x}</b><br>Budget: %{y:,.0f}<extra></extra>'
                            ))
                        
                        if not act_data.empty:
                            fig.add_trace(go.Bar(
                                x=act_data['Month'],
                                y=act_data['Value'],
                                name='Act',
                                marker_color='#FF8C00',
                                hovertemplate='<b>%{x}</b><br>Actual: %{y:,.0f}<extra></extra>'
                            ))
                        
                        fig.update_layout(
                            title={'text': f"Budget vs Actual Comparison - {table_name}", 'x': 0.5, 'xanchor': 'center', 'font': {'size': 16}},
                            xaxis_title="Month",
                            yaxis_title="Value",
                            font={'size': 12},
                            plot_bgcolor='white',
                            paper_bgcolor='white',
                            height=500,
                            margin={'l': 60, 'r': 60, 't': 80, 'b': 60},
                            showlegend=True,
                            barmode='group',
                            hovermode='x unified'
                        )
                        fig.update_xaxes(
                            title_font={'size': 14}, 
                            tickfont={'size': 12},
                            categoryorder='array',  # Use custom order
                            categoryarray=chart_data_agg['Month'].unique()  # Use our sorted months
                        )
                        fig.update_yaxes(title_font={'size': 14}, tickfont={'size': 12})
                        config = {'displayModeBar': True, 'displaylogo': False, 'modeBarButtonsToRemove': ['pan2d', 'lasso2d', 'select2d']}
                        st.plotly_chart(fig, use_container_width=True, config=config)
                        
                        ppt_type = 'bar' if visual_type == 'Bar Chart' else 'line'
                        chart_data_for_ppt = chart_data_agg.drop(columns=['MonthOnly', 'SortKey'])
                        x_col_for_ppt = "Month"
            
                    with st.expander("📊 View Data Table"):
                        st.dataframe(chart_data_for_ppt, use_container_width=True)

                    ppt_bytes = create_ppt_with_chart(
                        title=f"Budget vs Actual - {table_name} - {selected_sheet}",
                        chart_data=chart_data_for_ppt,
                        x_col=x_col_for_ppt,
                        y_col="Value",
                        chart_type=ppt_type,
                        selected_filter=selected_branch if is_branch_analysis else selected_product if is_product_analysis else None
                    )

                    st.download_button(
                        "⬇️ Download Budget vs Actual PPT",
                        ppt_bytes,
                        "budget_vs_actual.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_budget_vs_actual_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data_for_ppt

            def plot_monthly_comparison(tab, label, visual_type):
                with tab:
                    normalized_label = label.replace(",", "")
                    plot_cols = [col for col in table_df.columns 
                               if str(col).lower().replace(",", "").startswith(normalized_label.lower()) 
                               and 'ytd' not in str(col).lower()
                               and column_filter(col)]
                
                    if not plot_cols:
                        st.info(f"No matching columns found for '{label}'")
                        return None
                
                    chart_data = filtered_df[[first_col] + plot_cols].copy()
                
                    for col in plot_cols:
                        chart_data[col] = pd.to_numeric(chart_data[col].astype(str).str.replace(',', ''), 
                                                       errors='coerce')
                
                    chart_data = chart_data.melt(id_vars=first_col, 
                                              var_name="Month", 
                                              value_name=label)
                    chart_data = chart_data.dropna()
                
                    chart_data[label] = pd.to_numeric(chart_data[label], errors='coerce')
                
                    if chart_data.empty or not ensure_numeric_data(chart_data, label):
                        st.warning(f"No valid numeric data available for '{label}' after conversion.")
                        return None
                
                    chart_data['Month'] = chart_data['Month'].apply(extract_month_year)
                
                    month_order = {'Apr': 1, 'May': 2, 'Jun': 3, 'Jul': 4, 'Aug': 5, 'Sep': 6,
                                   'Oct': 7, 'Nov': 8, 'Dec': 9, 'Jan': 10, 'Feb': 11, 'Mar': 12}
                
                    def get_sort_key(month_str):
                        month_match = re.search(r'(\w{3,})[-–](\d{2})', month_str, re.IGNORECASE)
                        if month_match:
                            month, year = month_match.groups()
                            month_idx = month_order.get(month.capitalize(), 99)
                            year_int = int(year)
                            if month_idx >= 10:
                                fiscal_year = year_int - 1
                            else:
                                fiscal_year = year_int
                            return (fiscal_year, month_idx)
                        return (0, 99)
                
                    chart_data = chart_data.sort_values(by='Month', key=lambda x: x.map(get_sort_key))
                    chart_data = make_jsonly_serializable(chart_data)
                
                    color_override = '#FF8C00' if label == "Act" else None
                
                    display_visualization(tab, f"{label} by Month", chart_data, "Month", label, visual_type, color_override)
                    
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line' if visual_type == "Line Chart" else 'pie'
                    ppt_bytes = create_ppt_with_chart(
                        f"{label} Analysis - {table_name} - {selected_sheet}",
                        chart_data,
                        "Month",
                        label,
                        ppt_type,
                        color_override,
                        selected_filter=selected_branch if is_branch_analysis else selected_product if is_product_analysis else None
                    )
                    
                    st.download_button(
                        f"⬇️ Download {label} PPT",
                        ppt_bytes,
                        f"{label.lower().replace(' ', '_')}_analysis.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_{label.lower().replace(' ', '_')}_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data

            def plot_ytd_comparison(tab, pattern, label, visual_type):
                with tab:
                    ytd_cols = []
                    normalized_label = label.replace(",", "").lower()
                    
                    for col in table_df.columns:
                        col_str = str(col).lower().replace(",", "").replace("–", "-")
                        if normalized_label == 'gr':
                            if (re.search(r'gr-ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)', col_str, re.IGNORECASE) and
                                column_filter(col)):
                                ytd_cols.append(col)
                        elif normalized_label == 'ach':
                            if (re.search(r'ach-ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)', col_str, re.IGNORECASE) and
                                column_filter(col)):
                                ytd_cols.append(col)
                            elif (re.search(r'(?:ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\).*ach|ach.*ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))', col_str, re.IGNORECASE) and
                                  column_filter(col)):
                                ytd_cols.append(col)
                        else:
                            if (re.search(r'ytd.*\b' + normalized_label + r'\b|' + normalized_label + r'.*ytd', col_str, re.IGNORECASE) or
                                re.search(r'ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)\s*' + normalized_label, col_str, re.IGNORECASE)) and \
                               column_filter(col):
                                ytd_cols.append(col)
                    
                    if not ytd_cols:
                        st.warning(f"No YTD {label} columns found. Expected columns like '{label}-YTD-25-26-(Apr to Jun)'.")
                        return None
                    
                    clean_labels = []
                    for col in ytd_cols:
                        col_str = str(col)
                        year_match = re.search(r'(\d{2,4})\s*[-–]\s*(\d{2,4})\s*\((.*?)\)', col_str, re.IGNORECASE)
                        if year_match:
                            start_year, end_year, month_range = year_match.groups()
                            start_year = start_year[-2:] if len(start_year) > 2 else start_year
                            end_year = end_year[-2:] if len(end_year) > 2 else end_year
                            fiscal_year = f"{start_year}-{end_year}"
                            month_range_clean = re.sub(r'\s*to\s*', ' - ', month_range, flags=re.IGNORECASE)
                            if label.lower() == 'budget':
                                clean_label = f"Budget {fiscal_year} ({month_range_clean})"
                            elif label.lower() == 'ly':
                                clean_label = f"LY {fiscal_year} ({month_range_clean})"
                            else:
                                clean_label = f"{label} {fiscal_year} ({month_range_clean})"
                        else:
                            fiscal_year = "Unknown"
                            month_range_clean = "Apr - Jun"
                            if label.lower() == 'budget':
                                clean_label = f"Budget {fiscal_year} ({month_range_clean})"
                            elif label.lower() == 'ly':
                                clean_label = f"LY {fiscal_year} ({month_range_clean})"
                            else:
                                clean_label = f"{label} {fiscal_year} ({month_range_clean})"
                            st.warning(f"Could not parse year or month range in column '{col}'. Using default '{clean_label}'.")
                        clean_labels.append(clean_label)
                    
                    month_order = {'Apr':1, 'May':2, 'Jun':3, 'Jul':4, 'Aug':5, 'Sep':6,
                                   'Oct':7, 'Nov':8, 'Dec':9, 'Jan':10, 'Feb':11, 'Mar':12}
                    
                    def get_sort_key(col_name):
                        month_match = re.search(r'\((\w{3})', col_name, re.IGNORECASE)
                        return month_order.get(month_match.group(1).capitalize(), 0) if month_match else 0
                    
                    sorted_cols = [first_col] + sorted(clean_labels, key=get_sort_key)
                    comparison_data = filtered_df[[first_col] + ytd_cols].copy()
                    comparison_data.columns = [first_col] + clean_labels
                    comparison_data = comparison_data[sorted_cols]
                    
                    for col in clean_labels:
                        comparison_data[col] = pd.to_numeric(comparison_data[col].astype(str).str.replace(',', ''), errors='coerce')
                    
                    chart_data = comparison_data.melt(id_vars=first_col, 
                                                     var_name="Period", 
                                                     value_name=label)
                    chart_data = chart_data.dropna()
                    
                    if not ensure_numeric_data(chart_data, label):
                        st.warning(f"No numeric data available for YTD {label} comparisons")
                        return None
                    
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    color_override = '#FF8C00' if label == "Act" else None
                    
                    st.markdown(f"### {label} YTD Comparisons - {table_name}")
                    
                    fig, config = create_plotly_chart(
                        chart_data, 
                        "Period", 
                        label, 
                        visual_type.lower().replace(" chart", ""), 
                        f"{label} YTD Comparisons - {table_name}",
                        color_override
                    )
                    
                    # Force straight x-axis labels for YTD charts
                    fig.update_xaxes(tickangle=0)
                    
                    st.plotly_chart(fig, use_container_width=True, config=config)
                    
                    with st.expander("📊 View Data Table"):
                        st.dataframe(chart_data, use_container_width=True)
                    
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line' if visual_type == "Line Chart" else 'pie'
                    ppt_bytes = create_ppt_with_chart(
                        f"YTD {label} Analysis - {table_name} - {selected_sheet}",
                        chart_data,
                        "Period",
                        label,
                        ppt_type,
                        color_override,
                        selected_filter=selected_branch if is_branch_analysis else selected_product if is_product_analysis else None
                    )
                    
                    st.download_button(
                        f"⬇️ Download YTD {label} PPT",
                        ppt_bytes,
                        f"ytd_{label.lower().replace(' ', '_')}_analysis.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_ytd_{label.lower().replace(' ', '_')}_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data

            def plot_branch_performance(tab, visual_type):
                with tab:
                    if not is_branch_analysis:
                        st.info("This visualization is only available for region analysis sheets")
                        return None
                
                    ytd_act_col = None
                    for col in table_df.columns:
                        col_str = str(col).strip()
                        if col_str == "Act-YTD-25-26 (Apr to Mar)" or \
                           re.search(r'YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)\s*Act', col_str, re.IGNORECASE):
                            ytd_act_col = col
                            break
                
                    if ytd_act_col is None:
                        st.warning("Could not find YTD Act column for region performance analysis")
                        return None
                
                    first_col = table_df.columns[0]
                    regions_df = table_df[~table_df[first_col].str.contains('|'.join(BRANCH_EXCLUDE_TERMS), na=False, case=False)].copy()
                    # Remove rows starting with "GRAND TOTAL", "TOTAL SALES", or "OVERALL TOTAL"
                    regions_df = regions_df[~regions_df[first_col].str.contains(
                        r'^(?:TOTAL SALES|GRAND TOTAL|OVERALL TOTAL)', na=False, case=False, regex=True
                    )].copy()
                    regions_df = regions_df.dropna(subset=[first_col, ytd_act_col])
                
                    if regions_df.empty:
                        st.warning("No branch data available after filtering")
                        return None
                
                    regions_df[ytd_act_col] = pd.to_numeric(regions_df[ytd_act_col].astype(str).str.replace(',', ''), errors='coerce')
                    regions_df = regions_df.dropna(subset=[ytd_act_col])
                
                    if not ensure_numeric_data(regions_df, ytd_act_col):
                        st.warning("No numeric data available for region performance")
                        return None
                
                    regions_df = regions_df.sort_values(by=ytd_act_col, ascending=False)
                    
                    # Create a clean dataframe with only the columns we need
                    clean_regions_df = regions_df[[first_col, ytd_act_col]].copy()
                    clean_regions_df.columns = ['Branch', 'Performance']  # Standardize column names
                    
                    st.markdown(f"### Branch Performance Analysis - {table_name}")
                    
                    # Use the original data for display but return the clean data
                    display_visualization(tab, "Branch Performance", regions_df, first_col, ytd_act_col, visual_type)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        top_region = regions_df.iloc[0]
                        st.metric("Top Performer", top_region[first_col], f"{top_region[ytd_act_col]:,.0f}")
                    with col2:
                        total_performance = regions_df[ytd_act_col].sum()
                        st.metric("Total Performance", f"{total_performance:,.0f}")
                    with col3:
                        avg_performance = regions_df[ytd_act_col].mean()
                        st.metric("Average Performance", f"{avg_performance:,.0f}")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        st.markdown("#### 🏆 Top 5 Regions")
                        top_5 = regions_df.head(5)[[first_col, ytd_act_col]]
                        st.dataframe(top_5, use_container_width=True, hide_index=True)
                    
                    with col2:
                        st.markdown("#### 📉 Bottom 5 Regions")
                        bottom_5 = regions_df.tail(5)[[first_col, ytd_act_col]]
                        st.dataframe(bottom_5, use_container_width=True, hide_index=True)
                
                    clean_regions_df = make_jsonly_serializable(clean_regions_df)
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line' if visual_type == "Line Chart" else 'pie'
                    ppt_bytes = create_ppt_with_chart(
                        f"Branch Performance - {table_name} - {selected_sheet}",
                        clean_regions_df,
                        'Branch',
                        'Performance',
                        ppt_type,
                        selected_filter=selected_branch if selected_branch != "Select All" else None
                    )
                
                    st.download_button(
                        "⬇️ Download Region Performance PPT",
                        ppt_bytes,
                        "region_performance.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_region_performance_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return clean_regions_df

            def plot_branch_monthwise(tab, visual_type):
                with tab:
                    if not is_branch_analysis:
                        st.info("This visualization is only available for region analysis sheets")
                        return
                
                    act_cols = []
                    for col in table_df.columns:
                        col_str = str(col).lower()
                        for year in years:
                            if (re.search(rf'\bact\b.*(?:apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)[-\s]*{re.escape(year)}', col_str, re.IGNORECASE) and 'ytd' not in col_str):
                                act_cols.append(col)
                
                    if not act_cols:
                        st.warning(f"Could not find monthly Act columns for the selected years ({', '.join(years)})")
                        return
                
                    month_order = ['Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar']
                
                    def get_sort_key(col_name):
                        col_name = str(col_name).lower()
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', col_name, re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', col_name)
                        month_idx = month_order.index(month_match.group(1).capitalize()) if month_match and month_match.group(1).capitalize() in month_order else 99
                        year = int(year_match.group(1)) if year_match else 0
                        return (year, month_idx)
                
                    act_cols_sorted = sorted(act_cols, key=get_sort_key)
                
                    first_col = table_df.columns[0]
                    # Apply filtering here for Branch Monthwise visualization
                    regions_df = table_df[~table_df[first_col].str.contains('|'.join(BRANCH_EXCLUDE_TERMS), na=False, case=False, regex=True)].copy()
                    # Remove total rows for monthwise analysis
                    regions_df = regions_df[~regions_df[first_col].str.contains(
                        r'^(?:TOTAL SALES|GRAND TOTAL|OVERALL TOTAL)', na=False, case=False, regex=True
                    )].copy()
                    monthwise_data = regions_df[[first_col] + act_cols_sorted].copy()
                    
                    clean_col_names = []
                    for col in act_cols_sorted:
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', str(col), re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', str(col))
                        if month_match and year_match:
                            clean_col_names.append(f"{month_match.group(1).capitalize()}-{year_match.group(1)}")
                        else:
                            clean_col_names.append(str(col))
                    
                    monthwise_data.columns = [first_col] + clean_col_names
                    
                    for col in clean_col_names:
                        monthwise_data[col] = pd.to_numeric(monthwise_data[col].astype(str).str.replace(',', ''), errors='coerce')
                    
                    monthwise_data = monthwise_data.dropna()
                    
                    if monthwise_data.empty:
                        st.warning("No numeric data available for region monthwise performance after filtering")
                        return
                
                    st.write(f"### Branch Monthwise Performance ({', '.join(selected_years if selected_years else years)})")
                
                    chart_data = monthwise_data.melt(id_vars=first_col, 
                                                  var_name="Month", 
                                                  value_name="Value")
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    display_visualization(tab, "Branch Monthwise", chart_data, "Month", "Value", visual_type)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        total_by_month = monthwise_data[clean_col_names].sum()
                        best_month = total_by_month.idxmax()
                        st.metric("Best Month", best_month, f"{total_by_month[best_month]:,.0f}")
                    with col2:
                        avg_monthly = total_by_month.mean()
                        st.metric("Monthly Average", f"{avg_monthly:,.0f}")
                    with col3:
                        total_performance = total_by_month.sum()
                        st.metric("Total Performance", f"{total_performance:,.0f}")
                
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line'
                    ppt_bytes = create_ppt_with_chart(
                        f"Region Monthwise Performance - {selected_sheet}",
                        chart_data,
                        "Month",
                        "Value",
                        ppt_type,
                        selected_filter=selected_branch if selected_branch != "Select All" else None
                    )
                
                    st.download_button(
                        "⬇️ Download Region Monthwise PPT",
                        ppt_bytes,
                        "region_monthwise.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_region_monthwise_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data

            def plot_product_performance(tab, visual_type):
                with tab:
                    if not is_product_analysis:
                        st.info("This visualization is only available for product analysis sheets")
                        return None
                
                    ytd_act_col = None
                    for col in table_df.columns:
                        col_str = str(col).strip()
                        if col_str == "Act-YTD-25-26 (Apr to Mar)" or \
                           re.search(r'YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)\s*Act', col_str, re.IGNORECASE):
                            ytd_act_col = col
                            break
                
                    if ytd_act_col is None:
                        st.warning("Could not find YTD Act column for product performance analysis")
                        return None
                
                    first_col = table_df.columns[0]
                    exclude_terms = ['Total', 'TOTAL', 'Grand', 'GRAND', 'Total Sales']
                    products_df = table_df[~table_df[first_col].str.contains('|'.join(exclude_terms), na=False, case=False)].copy()
                    products_df = products_df.dropna(subset=[first_col, ytd_act_col])
                
                    if products_df.empty:
                        st.warning("No product data available after filtering")
                        return None
                
                    products_df[ytd_act_col] = pd.to_numeric(products_df[ytd_act_col].astype(str).str.replace(',', ''), errors='coerce')
                    products_df = products_df.dropna(subset=[ytd_act_col])
                
                    if not ensure_numeric_data(products_df, ytd_act_col):
                        st.warning("No numeric data available for product performance")
                        return None
                
                    products_df = products_df.sort_values(by=ytd_act_col, ascending=False)
                    
                    # Create a clean dataframe with only the columns we need
                    clean_products_df = products_df[[first_col, ytd_act_col]].copy()
                    clean_products_df.columns = ['Product', 'Performance']  # Standardize column names
                    
                    st.markdown("### Product Performance Analysis")
                    
                    # Use the original data for display but return the clean data
                    display_visualization(tab, "Product Performance", products_df, first_col, ytd_act_col, visual_type)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        top_product = products_df.iloc[0]
                        st.metric("Top Performer", top_product[first_col], f"{top_product[ytd_act_col]:,.0f}")
                    with col2:
                        total_performance = products_df[ytd_act_col].sum()
                        st.metric("Total Performance", f"{total_performance:,.0f}")
                    with col3:
                        avg_performance = products_df[ytd_act_col].mean()
                        st.metric("Average Performance", f"{avg_performance:,.0f}")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        st.markdown("#### 🏆 Top 5 Products")
                        top_5 = products_df.head(5)[[first_col, ytd_act_col]]
                        st.dataframe(top_5, use_container_width=True, hide_index=True)
                    
                    with col2:
                        st.markdown("#### 📉 Bottom 5 Products")
                        bottom_5 = products_df.tail(5)[[first_col, ytd_act_col]]
                        st.dataframe(bottom_5, use_container_width=True, hide_index=True)
                
                    clean_products_df = make_jsonly_serializable(clean_products_df)
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line' if visual_type == "Line Chart" else 'pie'
                    ppt_bytes = create_ppt_with_chart(
                        f"Product Performance - {selected_sheet}",
                        clean_products_df,
                        'Product',
                        'Performance',
                        ppt_type,
                        selected_filter=selected_product if selected_product != "Select All" else None
                    )
                
                    st.download_button(
                        "⬇️ Download Product Performance PPT",
                        ppt_bytes,
                        "product_performance.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_product_performance_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return clean_products_df

            def plot_product_monthwise(tab, visual_type):
                with tab:
                    if not is_product_analysis:
                        st.info("This visualization is only available for product analysis sheets")
                        return
                
                    act_cols = []
                    for col in table_df.columns:
                        col_str = str(col).lower()
                        for year in years:
                            if (re.search(rf'\bact\b.*(?:apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)[-\s]*{re.escape(year)}', col_str, re.IGNORECASE) and 'ytd' not in col_str):
                                act_cols.append(col)
                
                    if not act_cols:
                        st.warning(f"Could not find monthly Act columns for the selected years ({', '.join(years)})")
                        return
                
                    month_order = ['Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar']
                
                    def get_sort_key(col_name):
                        col_name = str(col_name).lower()
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', col_name, re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', col_name)
                        month_idx = month_order.index(month_match.group(1).capitalize()) if month_match and month_match.group(1).capitalize() in month_order else 99
                        year = int(year_match.group(1)) if year_match else 0
                        return (year, month_idx)
                
                    act_cols_sorted = sorted(act_cols, key=get_sort_key)
                
                    first_col = table_df.columns[0]
                    # Apply filtering here for Product Monthwise visualization
                    exclude_terms = ['Total', 'TOTAL', 'Grand', 'GRAND', 'Total Sales']
                    products_df = table_df[~table_df[first_col].str.contains('|'.join(exclude_terms), na=False, case=False, regex=True)].copy()
                    # Remove total rows for monthwise analysis
                    products_df = products_df[~products_df[first_col].str.contains(
                        r'^(?:TOTAL SALES|GRAND TOTAL|OVERALL TOTAL)', na=False, case=False, regex=True
                    )].copy()
                    monthwise_data = products_df[[first_col] + act_cols_sorted].copy()
                    
                    clean_col_names = []
                    for col in act_cols_sorted:
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', str(col), re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', str(col))
                        if month_match and year_match:
                            clean_col_names.append(f"{month_match.group(1).capitalize()}-{year_match.group(1)}")
                        else:
                            clean_col_names.append(str(col))
                    
                    monthwise_data.columns = [first_col] + clean_col_names
                    
                    for col in clean_col_names:
                        monthwise_data[col] = pd.to_numeric(monthwise_data[col].astype(str).str.replace(',', ''), errors='coerce')
                    
                    monthwise_data = monthwise_data.dropna()
                    
                    if monthwise_data.empty:
                        st.warning("No numeric data available for product monthwise performance after filtering")
                        return
                
                    st.write(f"### Product Monthwise Performance ({', '.join(selected_years if selected_years else years)})")
                
                    chart_data = monthwise_data.melt(id_vars=first_col, 
                                                  var_name="Month", 
                                                  value_name="Value")
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    display_visualization(tab, "Product Monthwise", chart_data, "Month", "Value", visual_type)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        total_by_month = monthwise_data[clean_col_names].sum()
                        best_month = total_by_month.idxmax()
                        st.metric("Best Month", best_month, f"{total_by_month[best_month]:,.0f}")
                    with col2:
                        avg_monthly = total_by_month.mean()
                        st.metric("Monthly Average", f"{avg_monthly:,.0f}")
                    with col3:
                        total_performance = total_by_month.sum()
                        st.metric("Total Performance", f"{total_performance:,.0f}")
                
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line'
                    ppt_bytes = create_ppt_with_chart(
                        f"Product Monthwise Performance - {selected_sheet}",
                        chart_data,
                        "Month",
                        "Value",
                        ppt_type,
                        selected_filter=selected_product if selected_product != "Select All" else None
                    )
                
                    st.download_button(
                        "⬇️ Download Product Monthwise PPT",
                        ppt_bytes,
                        "product_monthwise.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_product_monthwise_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data

            # Plot visualizations in respective tabs
            budget_vs_actual_data = plot_budget_vs_actual(tabs_dict["Budget vs Actual"], visual_type)
            budget_data = plot_monthly_comparison(tabs_dict["Budget"], "Budget", visual_type)
            ly_data = plot_monthly_comparison(tabs_dict["LY"], "LY", visual_type)
            act_data = plot_monthly_comparison(tabs_dict["Act"], "Act", visual_type)
            gr_data = plot_monthly_comparison(tabs_dict["Gr"], "Gr", visual_type)
            ach_data = plot_monthly_comparison(tabs_dict["Ach"], "Ach", visual_type)
            ytd_budget_data = plot_ytd_comparison(tabs_dict["YTD Budget"], r'\bBudget\b.*YTD', "Budget", visual_type)
            ytd_ly_data = plot_ytd_comparison(tabs_dict["YTD LY"], r'\bLY\b.*YTD', "LY", visual_type)
            ytd_act_data = plot_ytd_comparison(tabs_dict["YTD Act"], r'\bAct\b.*YTD', "Act", visual_type)
            ytd_gr_data = plot_ytd_comparison(tabs_dict["YTD Gr"], r'\bGr\b.*YTD', "Gr", visual_type)
            ytd_ach_data = plot_ytd_comparison(tabs_dict["YTD Ach"], r'\bAch\b.*YTD', "Ach", visual_type)
            branch_performance_data = plot_branch_performance(tabs_dict["Branch Performance"], visual_type)
            branch_monthwise_data = plot_branch_monthwise(tabs_dict["Branch Monthwise"], visual_type)
            product_performance_data = plot_product_performance(tabs_dict["Product Performance"], visual_type)
            product_monthwise_data = plot_product_monthwise(tabs_dict["Product Monthwise"], visual_type)

            # Generate master PPT for all visualizations - CHARTS ONLY, NO TABLES
            all_data = [
                ("Budget vs Actual", budget_vs_actual_data),
                ("Budget", budget_data),
                ("LY", ly_data),
                ("Act", act_data),
                ("Gr", gr_data),
                ("Ach", ach_data),
                ("YTD Budget", ytd_budget_data),
                ("YTD LY", ytd_ly_data),
                ("YTD Act", ytd_act_data),
                ("YTD Gr", ytd_gr_data),
                ("YTD Ach", ytd_ach_data),
                ("Branch Performance", branch_performance_data),
                ("Branch Monthwise", branch_monthwise_data),
                ("Product Performance", product_performance_data),
                ("Product Monthwise", product_monthwise_data)
            ]

            if any(data is not None for _, data in all_data):
                st.sidebar.markdown("---")
                st.sidebar.subheader("📊 Download All Visuals")
                
                # Get the selected filter for the master PPT title
                selected_filter = None
                if is_branch_analysis and selected_branch != "Select All":
                    selected_filter = selected_branch
                elif is_product_analysis and selected_product != "Select All":
                    selected_filter = selected_product
                
                if st.sidebar.button("🔄 Generate Master PPT", help="Create PPT with all charts (NO tables, straight x-axis labels)"):
                    with st.spinner("Generating PowerPoint with all charts..."):
                        try:
                            master_ppt_bytes = create_master_ppt_with_matplotlib(
                                all_data, 
                                table_name, 
                                selected_sheet, 
                                visual_type,
                                selected_filter
                            )
                            
                            if master_ppt_bytes:
                                st.sidebar.success("✅ Master PPT generated successfully!")
                                st.sidebar.download_button(
                                    "⬇️ Download Master PPT (Charts Only)",
                                    master_ppt_bytes,
                                    f"charts_only_{selected_sheet}.pptx",
                                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key=f"download_master_ppt_{selected_sheet}_{sheet_index}"
                                )
                            else:
                                st.sidebar.error("❌ Failed to generate master PPT")
                                
                        except Exception as e:
                            st.sidebar.error(f"❌ Error generating master PPT: {str(e)}")
                
                # Alternative individual download approach
                st.sidebar.markdown("#### 📄 Individual Chart Downloads")
                st.sidebar.info("💡 All PPT downloads contain ONLY clean charts (no values, no tables, straight x-axis labels)")

            optimize_memory()

else:
    st.info("Please upload an Excel file to begin analysis.")
    
optimize_memory()
